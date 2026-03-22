from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Callable

from fastapi import HTTPException, status

from replace_pdf_bank_info import ensure_pymupdf


PPTX_CONTENT_TYPES = {
    ".pptx",
}


def ensure_python_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Cm  # type: ignore
    except ImportError as exc:  # pragma: no cover - depends on deployment env
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="python-pptx is required to process PowerPoint files",
        ) from exc
    return Presentation, Cm


def validate_pptx_filename(filename: str) -> None:
    suffix = Path(filename or "").suffix.lower()
    if suffix not in PPTX_CONTENT_TYPES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Only .pptx files are supported",
        )


def count_presentation_slides(file_path: Path) -> int:
    Presentation, _ = ensure_python_pptx()
    presentation = Presentation(str(file_path))
    return len(presentation.slides)


def _find_libreoffice_binary() -> str:
    for candidate in ("libreoffice", "soffice"):
        binary = shutil.which(candidate)
        if binary:
            return binary
    raise HTTPException(
        status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
        detail="LibreOffice is required to render PowerPoint files on the server",
    )


def _render_pptx_to_pdf(input_path: Path, work_dir: Path) -> Path:
    binary = _find_libreoffice_binary()
    profile_dir = work_dir / "lo-profile"
    profile_dir.mkdir(parents=True, exist_ok=True)
    output_dir = work_dir / "pdf"
    output_dir.mkdir(parents=True, exist_ok=True)

    command = [
        binary,
        "--headless",
        f"-env:UserInstallation=file://{profile_dir.as_posix()}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_path),
    ]

    result = subprocess.run(
        command,
        capture_output=True,
        text=True,
        timeout=240,
        check=False,
    )
    if result.returncode != 0:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"LibreOffice failed to convert PPTX: {(result.stderr or result.stdout).strip()}",
        )

    pdf_path = output_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Rendered PDF was not created by LibreOffice",
        )
    return pdf_path


def _fit_size(img_width: int, img_height: int, box_width: int, box_height: int) -> tuple[int, int]:
    img_ratio = img_width / img_height
    box_ratio = box_width / box_height
    if img_ratio >= box_ratio:
        width = box_width
        height = int(width / img_ratio)
    else:
        height = box_height
        width = int(height * img_ratio)
    return width, height


def convert_presentation_to_three_panel(
    input_path: Path,
    output_path: Path,
    progress_callback: Callable[[int, int], None] | None = None,
) -> int:
    Presentation, Cm = ensure_python_pptx()
    fitz = ensure_pymupdf()

    with tempfile.TemporaryDirectory(prefix="ppt_three_panel_") as temp_dir_raw:
        temp_dir = Path(temp_dir_raw)
        rendered_pdf = _render_pptx_to_pdf(input_path, temp_dir)

        source_pdf = fitz.open(rendered_pdf)
        try:
            total_slides = source_pdf.page_count
            presentation = Presentation()
            presentation.slide_width = Cm(101.6)
            presentation.slide_height = Cm(19.05)
            blank_layout = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[-1]

            # Remove the default first slide if present.
            if presentation.slides:
                slide_id_list = presentation.slides._sldIdLst  # type: ignore[attr-defined]
                if len(slide_id_list):
                    slide_id_list.remove(slide_id_list[0])

            panel_width = int(presentation.slide_width / 3)
            slide_height = int(presentation.slide_height)
            render_dir = temp_dir / "rendered"
            render_dir.mkdir(parents=True, exist_ok=True)

            for page_index in range(total_slides):
                page = source_pdf.load_page(page_index)
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                image_path = render_dir / f"slide_{page_index + 1}.png"
                pixmap.save(image_path)

                slide = presentation.slides.add_slide(blank_layout)
                target_width, target_height = _fit_size(pixmap.width, pixmap.height, panel_width, slide_height)
                top = int((slide_height - target_height) / 2)

                for panel_index in range(3):
                    left = int(panel_index * panel_width + (panel_width - target_width) / 2)
                    slide.shapes.add_picture(str(image_path), left, top, width=target_width, height=target_height)

                if progress_callback is not None:
                    progress_callback(page_index + 1, total_slides)

            output_path.parent.mkdir(parents=True, exist_ok=True)
            presentation.save(str(output_path))
        finally:
            source_pdf.close()

    if progress_callback is not None:
        progress_callback(total_slides, total_slides)
    return total_slides
